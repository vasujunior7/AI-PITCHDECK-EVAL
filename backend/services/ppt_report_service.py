"""
📊 PowerPoint Report Service - Milestone 9 Extension
Generates professional PowerPoint reports with analysis results
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from datetime import datetime
from pathlib import Path
import io
from typing import Dict, List
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import matplotlib.pyplot as plt
import numpy as np
from io import BytesIO
from core.logging import logger


class PPTReportGenerator:
    """Generate comprehensive PowerPoint analysis reports"""
    
    # Color scheme (matching academic theme)
    COLOR_PRIMARY = RGBColor(41, 128, 185)  # Blue
    COLOR_SUCCESS = RGBColor(39, 174, 96)   # Green
    COLOR_WARNING = RGBColor(230, 126, 34)  # Orange
    COLOR_DANGER = RGBColor(231, 76, 60)    # Red
    COLOR_NEUTRAL = RGBColor(149, 165, 166) # Gray
    COLOR_DARK = RGBColor(44, 62, 80)       # Dark Blue
    COLOR_LIGHT = RGBColor(236, 240, 241)   # Light Gray
    
    def __init__(self):
        """Initialize report generator"""
        self.prs = None
    
    def _get_grade_color(self, grade: str) -> RGBColor:
        """Get color based on grade (updated for lenient scale)"""
        grade_letter = grade[0] if grade else 'F'  # Extract first letter
        
        if grade_letter == 'A':
            return self.COLOR_SUCCESS  # Green
        elif grade_letter == 'B':
            return self.COLOR_PRIMARY  # Blue
        elif grade_letter == 'C':
            return self.COLOR_WARNING  # Orange
        elif grade_letter == 'D' or grade_letter == 'F':
            return self.COLOR_DANGER  # Red
        else:
            return self.COLOR_NEUTRAL
    
    def _get_score_color(self, score: float) -> RGBColor:
        """Get color based on score (updated for lenient scale)"""
        if score >= 75:
            return self.COLOR_SUCCESS  # Green - A
        elif score >= 55:
            return self.COLOR_PRIMARY  # Blue - B
        elif score >= 40:
            return self.COLOR_WARNING  # Orange - C
        elif score >= 25:
            return self.COLOR_DANGER  # Red - D
        else:
            return self.COLOR_DANGER  # Red - F
    
    def _add_title_slide(self, title: str, subtitle: str):
        """Add title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Add background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLOR_PRIMARY
        
        # Title
        left = Inches(0.5)
        top = Inches(2.5)
        width = Inches(9)
        height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(left, top + Inches(1.2), width, Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    
    def _add_section_slide(self, title: str):
        """Add section divider slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLOR_DARK
        
        # Title
        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(1.5)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
    
    def _add_content_slide(self, title: str, bullet_points: list):
        """Add content slide with bullet points"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[1])  # Title and Content
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.COLOR_DARK
        
        # Content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(18)
            p.space_before = Pt(12)
    
    def _add_overall_score_slide(self, overall_score: float, grade: str, 
                                  presentation_name: str, total_slides: int):
        """Add overall score summary slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Overall Score Summary"
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLOR_DARK
        
        # Score circle (large centered score)
        left = Inches(3.5)
        top = Inches(2)
        width = Inches(3)
        height = Inches(2)
        
        score_box = slide.shapes.add_textbox(left, top, width, height)
        score_frame = score_box.text_frame
        score_frame.text = f"{overall_score:.1f}/100"
        score_para = score_frame.paragraphs[0]
        score_para.alignment = PP_ALIGN.CENTER
        score_para.font.size = Pt(60)
        score_para.font.bold = True
        score_para.font.color.rgb = self._get_score_color(overall_score)
        
        # Grade
        grade_box = slide.shapes.add_textbox(left, top + Inches(1.5), width, Inches(0.8))
        grade_frame = grade_box.text_frame
        grade_frame.text = f"Grade: {grade}"
        grade_para = grade_frame.paragraphs[0]
        grade_para.alignment = PP_ALIGN.CENTER
        grade_para.font.size = Pt(32)
        grade_para.font.bold = True
        grade_para.font.color.rgb = self._get_grade_color(grade)
        
        # Metadata
        meta_text = f"{presentation_name} • {total_slides} Slides • {datetime.now().strftime('%B %d, %Y')}"
        meta_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
        meta_frame = meta_box.text_frame
        meta_frame.text = meta_text
        meta_para = meta_frame.paragraphs[0]
        meta_para.alignment = PP_ALIGN.CENTER
        meta_para.font.size = Pt(14)
        meta_para.font.color.rgb = self.COLOR_NEUTRAL
    
    def _add_component_scores_slide(self, components: dict):
        """Add component scores chart slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Component Scores Breakdown"
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLOR_DARK
        
        # Create bar chart
        chart_data = CategoryChartData()
        chart_data.categories = [comp['name'] for comp in components]
        chart_data.add_series('Score', [comp['score'] for comp in components])
        
        x, y, cx, cy = Inches(1.5), Inches(1.8), Inches(7), Inches(4)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        
        # Style chart
        chart.has_legend = False
        chart.has_title = False
        
        # Color bars based on score
        for idx, comp in enumerate(components):
            point = chart.series[0].points[idx]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = self._get_score_color(comp['score'])
    
    def _add_per_slide_scores_slide(self, slide_scores: list):
        """Add per-slide score distribution (line chart)"""
        if not slide_scores or len(slide_scores) == 0:
            return  # Skip if no scores available
            
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Per-Slide Score Distribution"
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLOR_DARK
        
        # Create line chart
        chart_data = CategoryChartData()
        chart_data.categories = [f"Slide {i+1}" for i in range(len(slide_scores))]
        chart_data.add_series('Score', slide_scores)
        
        x, y, cx, cy = Inches(1.5), Inches(1.8), Inches(7), Inches(4)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data
        ).chart
        
        chart.has_legend = False
        chart.has_title = False
        
        # Color line
        chart.series[0].format.line.color.rgb = self.COLOR_PRIMARY
        chart.series[0].format.line.width = Pt(3)
    
    def _add_strengths_slide(self, strengths: list):
        """Add strengths slide with checkmarks"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "✓ Key Strengths"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLOR_SUCCESS
        
        # Add strengths as bullet points
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for idx, strength in enumerate(strengths[:6]):  # Limit to 6
            p = content_frame.add_paragraph() if idx > 0 else content_frame.paragraphs[0]
            p.text = f"✓ {strength}"
            p.font.size = Pt(20)
            p.space_before = Pt(16)
            p.font.color.rgb = self.COLOR_DARK
    
    def _add_improvements_slide(self, improvements: list):
        """Add areas for improvement slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "⚠ Areas for Improvement"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLOR_WARNING
        
        # Add improvements as bullet points
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        for idx, improvement in enumerate(improvements[:6]):  # Limit to 6
            p = content_frame.add_paragraph() if idx > 0 else content_frame.paragraphs[0]
            p.text = f"→ {improvement}"
            p.font.size = Pt(20)
            p.space_before = Pt(16)
            p.font.color.rgb = self.COLOR_DARK
    
    def _add_recommendations_slide(self, recommendations: list):
        """Add top recommendations slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        
        # Title (clear that this is actionable recommendations / action plan)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "💡 Top Recommendations — Action Plan"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLOR_PRIMARY
        
        # Add top 5 recommendations
        y_pos = 1.8
        for idx, rec in enumerate(recommendations[:5]):
            # Number badge
            badge_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(0.4), Inches(0.4))
            badge_frame = badge_box.text_frame
            badge_frame.text = str(idx + 1)
            badge_para = badge_frame.paragraphs[0]
            badge_para.alignment = PP_ALIGN.CENTER
            badge_para.font.size = Pt(16)
            badge_para.font.bold = True
            badge_para.font.color.rgb = RGBColor(255, 255, 255)
            
            # Add colored background to badge
            shape = badge_box
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = self.COLOR_PRIMARY
            
            # Recommendation text - accept multiple possible keys from recommendation dict
            rec_box = slide.shapes.add_textbox(Inches(1.6), Inches(y_pos), Inches(7.4), Inches(0.7))
            rec_frame = rec_box.text_frame
            rec_frame.word_wrap = True
            # Normalise candidate fields
            rec_text = None
            for key in ('recommendation', 'description', 'message', 'issue'):
                if isinstance(rec, dict) and rec.get(key):
                    rec_text = rec.get(key)
                    break
            if not rec_text:
                # Fall back to category + short summary
                if isinstance(rec, dict):
                    cat = rec.get('category', '')
                    rec_text = f"{cat}: {rec.get('recommendation', rec.get('description', 'No description'))}"
                else:
                    rec_text = str(rec)

            rec_frame.text = rec_text
            rec_para = rec_frame.paragraphs[0]
            rec_para.font.size = Pt(16)
            rec_para.font.color.rgb = self.COLOR_DARK
            
            y_pos += 0.85
    
    def _add_feature_analysis_slide(self, features: dict):
        """Add feature analysis slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Feature Analysis Metrics"
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLOR_DARK
        
        # Helper to safely extract float values
        def safe_float(value, default=0):
            if isinstance(value, dict):
                for key in ["overall_density", "overall_quality", "average_level", "redundancy_score", "score"]:
                    if key in value:
                        try:
                            return float(value[key])
                        except (ValueError, TypeError):
                            continue
                return default
            try:
                return float(value)
            except (ValueError, TypeError):
                return default
        
        # Create 2x2 grid of metrics
        density = safe_float(features.get('semantic_density', 0))
        redundancy = safe_float(features.get('redundancy_score', 0))
        layout = safe_float(features.get('layout_quality', 0))
        blooms = safe_float(features.get('avg_blooms_level', 0))
        
        metrics = [
            ("Semantic Density", f"{density:.3f}", "Word relationships"),
            ("Redundancy", f"{redundancy:.3f}", "Content repetition"),
            ("Layout Quality", f"{layout:.1f}/100", "Visual design"),
            ("Bloom's Level", f"{blooms:.2f}/6", "Cognitive depth")
        ]
        
        x_positions = [1, 5.5]
        y_positions = [2, 4]
        
        for idx, (metric_name, metric_value, metric_desc) in enumerate(metrics):
            x = Inches(x_positions[idx % 2])
            y = Inches(y_positions[idx // 2])
            
            # Metric box
            box = slide.shapes.add_textbox(x, y, Inches(3.5), Inches(1.5))
            frame = box.text_frame
            
            # Metric name
            p1 = frame.paragraphs[0]
            p1.text = metric_name
            p1.font.size = Pt(18)
            p1.font.bold = True
            p1.font.color.rgb = self.COLOR_DARK
            p1.alignment = PP_ALIGN.CENTER
            
            # Metric value
            p2 = frame.add_paragraph()
            p2.text = metric_value
            p2.font.size = Pt(28)
            p2.font.bold = True
            p2.font.color.rgb = self.COLOR_PRIMARY
            p2.alignment = PP_ALIGN.CENTER
            
            # Metric description
            p3 = frame.add_paragraph()
            p3.text = metric_desc
            p3.font.size = Pt(12)
            p3.font.color.rgb = self.COLOR_NEUTRAL
            p3.alignment = PP_ALIGN.CENTER
    
    def _add_milestone3_evaluation_slide(self, evaluation: Dict):
        """Add a slide for a single Milestone 3 evaluation question"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        
        question_name = evaluation.get("question_name", "Unknown")
        rating_level = evaluation.get("rating_level", "N/A")
        justification = evaluation.get("justification", "")
        improvement = evaluation.get("improvement_suggestion", "")
        score = evaluation.get("score", 0)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = f"Question {evaluation.get('question_id', '?')}: {question_name}"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLOR_DARK
        
        # Rating badge
        rating_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(3), Inches(0.6))
        rating_frame = rating_box.text_frame
        rating_frame.text = f"Rating: {rating_level}"
        rating_para = rating_frame.paragraphs[0]
        rating_para.font.size = Pt(20)
        rating_para.font.bold = True
        rating_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Color rating badge based on score
        rating_shape = rating_box
        fill = rating_shape.fill
        fill.solid()
        fill.fore_color.rgb = self._get_score_color(score)
        
        # Score
        score_box = slide.shapes.add_textbox(Inches(4), Inches(1.2), Inches(2), Inches(0.6))
        score_frame = score_box.text_frame
        score_frame.text = f"Score: {score:.1f}/100"
        score_para = score_frame.paragraphs[0]
        score_para.font.size = Pt(18)
        score_para.font.bold = True
        score_para.font.color.rgb = self._get_score_color(score)
        
        # Justification
        just_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
        just_frame = just_box.text_frame
        just_frame.word_wrap = True
        just_frame.text = f"Justification:\n{justification[:400]}"
        just_para = just_frame.paragraphs[0]
        just_para.font.size = Pt(14)
        just_para.font.color.rgb = self.COLOR_DARK
        
        # Improvement suggestion
        improve_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1.5))
        improve_frame = improve_box.text_frame
        improve_frame.word_wrap = True
        improve_frame.text = f"Improvement Suggestion:\n{improvement[:400]}"
        improve_para = improve_frame.paragraphs[0]
        improve_para.font.size = Pt(14)
        improve_para.font.color.rgb = self.COLOR_PRIMARY
    
    def _add_milestone3_summary_slide(self, milestone3_data: Dict):
        """Add Milestone 3 overall summary slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        
        overall_score = milestone3_data.get('overall_score', 0)
        grade = milestone3_data.get('grade', 'N/A')
        category_scores = milestone3_data.get('category_scores', {})
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Milestone 3: Venture Assessment Summary"
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLOR_DARK
        
        # Overall score
        score_box = slide.shapes.add_textbox(Inches(3), Inches(1.5), Inches(4), Inches(1.5))
        score_frame = score_box.text_frame
        score_frame.text = f"{overall_score:.1f}/100"
        score_para = score_frame.paragraphs[0]
        score_para.alignment = PP_ALIGN.CENTER
        score_para.font.size = Pt(60)
        score_para.font.bold = True
        score_para.font.color.rgb = self._get_score_color(overall_score)
        
        # Grade
        grade_box = slide.shapes.add_textbox(Inches(3), Inches(2.8), Inches(4), Inches(0.8))
        grade_frame = grade_box.text_frame
        grade_frame.text = f"Grade: {grade}"
        grade_para = grade_frame.paragraphs[0]
        grade_para.alignment = PP_ALIGN.CENTER
        grade_para.font.size = Pt(32)
        grade_para.font.bold = True
        grade_para.font.color.rgb = self._get_grade_color(grade)
        
        # Category scores
        y_pos = 4.0
        for category, score in category_scores.items():
            cat_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(8), Inches(0.4))
            cat_frame = cat_box.text_frame
            cat_frame.text = f"{category}: {score:.1f}/100"
            cat_para = cat_frame.paragraphs[0]
            cat_para.font.size = Pt(16)
            cat_para.font.color.rgb = self.COLOR_DARK
            y_pos += 0.5
    
    def _generate_spider_chart(self, data: Dict, title: str, labels: List[str], output_path: str):
        """Generate a spider/radar chart and save as image"""
        try:
            # Prepare data
            values = [data.get(label, 0) for label in labels]
            
            # Number of variables
            N = len(labels)
            
            # Compute angle for each axis
            angles = [n / float(N) * 2 * np.pi for n in range(N)]
            angles += angles[:1]  # Complete the circle
            
            # Add first value to end to close the plot
            values += values[:1]
            
            # Create figure
            fig, ax = plt.subplots(figsize=(8, 8), subplot_kw=dict(projection='polar'))
            
            # Convert RGBColor to matplotlib color (normalize 0-255 to 0-1)
            # RGBColor: use rgb property or direct access
            try:
                r, g, b = self.COLOR_PRIMARY.r, self.COLOR_PRIMARY.g, self.COLOR_PRIMARY.b
            except:
                # Fallback: use blue color directly
                r, g, b = 41, 128, 185
            color = (r / 255.0, g / 255.0, b / 255.0)
            
            # Plot
            ax.plot(angles, values, 'o-', linewidth=2, color=color)
            ax.fill(angles, values, alpha=0.25, color=color)
            
            # Add labels
            ax.set_xticks(angles[:-1])
            ax.set_xticklabels(labels, fontsize=10)
            ax.set_ylim(0, 100)
            ax.set_yticks([20, 40, 60, 80, 100])
            ax.set_yticklabels(['20', '40', '60', '80', '100'], fontsize=8)
            ax.grid(True)
            
            # Title
            plt.title(title, size=14, fontweight='bold', pad=20)
            
            # Save
            plt.tight_layout()
            plt.savefig(output_path, dpi=150, bbox_inches='tight')
            plt.close()
            
            return output_path
        except Exception as e:
            logger.error(f"Error generating spider chart: {e}")
            return None
    
    def _add_spider_chart_slide(self, chart_image_path: str, title: str):
        """Add a slide with a spider chart image"""
        if not chart_image_path or not Path(chart_image_path).exists():
            return
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLOR_DARK
        
        # Add image
        slide.shapes.add_picture(chart_image_path, Inches(1), Inches(1.2), Inches(8), Inches(5.5))
    
    def _add_conclusion_slide(self, grade: str, key_message: str):
        """Add conclusion slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank
        
        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._get_grade_color(grade)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Analysis Complete"
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Key message
        msg_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(2))
        msg_frame = msg_box.text_frame
        msg_frame.word_wrap = True
        msg_frame.text = key_message
        msg_para = msg_frame.paragraphs[0]
        msg_para.alignment = PP_ALIGN.CENTER
        msg_para.font.size = Pt(20)
        msg_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Footer
        footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        footer_frame = footer_box.text_frame
        footer_frame.text = "AI-Powered Presentation Analyzer • Academic Grading System"
        footer_para = footer_frame.paragraphs[0]
        footer_para.alignment = PP_ALIGN.CENTER
        footer_para.font.size = Pt(12)
        footer_para.font.color.rgb = RGBColor(255, 255, 255)
    
    def generate_report(self, analysis_data: dict, output_path: str = None) -> str:
        """
        Generate comprehensive PowerPoint report
        
        Args:
            analysis_data: Complete analysis results
            output_path: Optional custom output path
            
        Returns:
            Path to generated PPTX file
        """
        # Create presentation
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Extract data
        presentation_name = analysis_data.get('presentation_name', 'Unknown Presentation')
        overall_score = analysis_data.get('overall_score', 0)
        grade = analysis_data.get('grade', 'N/A')
        components = analysis_data.get('component_scores', [])
        per_slide = analysis_data.get('slide_scores', [])  # Fixed: was 'per_slide_scores'
        features = analysis_data.get('feature_analysis', {})
        ai_eval = analysis_data.get('ai_evaluation', {})
        # Normalize recommendations: accept dict (with 'recommendations'), list, or other iterables
        recommendations = analysis_data.get('recommendations', [])
        try:
            if isinstance(recommendations, dict):
                # If it's the full recommendations dict produced for PDF, extract the list
                recommendations = recommendations.get('recommendations', []) or []
            elif isinstance(recommendations, slice):
                # Defensive: if a slice object was passed accidentally
                recommendations = []
            elif not isinstance(recommendations, list):
                # Try to coerce other iterables to list
                try:
                    recommendations = list(recommendations)
                except Exception:
                    recommendations = []
        except Exception:
            # Fallback to empty list on any unexpected issue
            recommendations = []
        
        # 1. Title Slide
        self._add_title_slide(
            "Presentation Analysis Report",
            f"Venture Assessment (Milestone 3) • {datetime.now().strftime('%B %Y')}"
        )
        
        # 2. Milestone 3 Evaluation Section (90% weight)
        milestone3_data = analysis_data.get('milestone3_evaluation', {})
        if milestone3_data:
            self._add_section_slide("Milestone 3: Venture Assessment Evaluation")
            
            # Milestone 3 Summary
            self._add_milestone3_summary_slide(milestone3_data)
            
            # Individual question evaluations (18 questions)
            evaluations = milestone3_data.get('evaluations', [])
            for eval_data in evaluations:
                self._add_milestone3_evaluation_slide(eval_data)
            
            # Spider Charts Section
            self._add_section_slide("Milestone 3: Visual Analysis")
            
            # Generate spider charts
            chart_dir = Path("outputs/reports/charts")
            chart_dir.mkdir(parents=True, exist_ok=True)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            
            # Chart 1: All 18 Milestone 3 questions
            question_labels = [f"Q{eval['question_id']}" for eval in evaluations]
            question_scores = {f"Q{eval['question_id']}": eval.get('score', 0) for eval in evaluations}
            chart1_path = chart_dir / f"milestone3_all_questions_{timestamp}.png"
            self._generate_spider_chart(question_scores, "Milestone 3: All 18 Questions", question_labels, str(chart1_path))
            if chart1_path.exists():
                self._add_spider_chart_slide(str(chart1_path), "Milestone 3: All 18 Questions")
            
            # Chart 2: Category scores
            category_scores = milestone3_data.get('category_scores', {})
            chart2_path = chart_dir / f"milestone3_categories_{timestamp}.png"
            self._generate_spider_chart(category_scores, "Milestone 3: Category Analysis", list(category_scores.keys()), str(chart2_path))
            if chart2_path.exists():
                self._add_spider_chart_slide(str(chart2_path), "Milestone 3: Category Analysis")
            
            # Chart 3: Comparison (Milestone 3 vs Old PPT scoring)
            old_ppt_score = analysis_data.get('overall_score', 0)  # This is already combined
            # Calculate old PPT score before combination (approximate)
            milestone3_score = milestone3_data.get('overall_score', 0)
            old_ppt_only = (overall_score - (milestone3_score * 0.90)) / 0.10 if milestone3_score > 0 else overall_score
            comparison_data = {
                "Milestone 3": milestone3_score,
                "Old PPT Scoring": old_ppt_only
            }
            chart3_path = chart_dir / f"milestone3_comparison_{timestamp}.png"
            self._generate_spider_chart(comparison_data, "Scoring Comparison", list(comparison_data.keys()), str(chart3_path))
            if chart3_path.exists():
                self._add_spider_chart_slide(str(chart3_path), "Scoring Comparison: Milestone 3 vs Old PPT")
        
        # 3. Overall Score Summary (Combined)
        self._add_section_slide("Overall Analysis")
        self._add_overall_score_slide(
            overall_score,
            grade,
            presentation_name,
            len(per_slide)
        )
        
        # 4. Section: Detailed Analysis (Old PPT Scoring - 10% weight)
        self._add_section_slide("Detailed Analysis (10% Weight)")
        
        # 4. Component Scores
        component_list = []
        for comp in components:
            component_list.append({
                'name': comp.get('component', 'Unknown'),
                'score': comp.get('score', 0)
            })
        self._add_component_scores_slide(component_list)
        
        # 5. Per-Slide Scores
        slide_scores = [s.get('score', 0) for s in per_slide]
        self._add_per_slide_scores_slide(slide_scores)
        
        # 6. Feature Analysis
        self._add_feature_analysis_slide(features)
        
        # 7. Section: Evaluation Results
        self._add_section_slide("Evaluation Results")
        
        # 8. Strengths
        strengths = ai_eval.get('strengths', [])
        if strengths:
            self._add_strengths_slide(strengths)
        
        # 9. Areas for Improvement
        improvements = ai_eval.get('areas_for_improvement', [])
        if improvements:
            self._add_improvements_slide(improvements)
        
        # 10. Section: Recommendations
        self._add_section_slide("Action Plan")
        
        # 11. Top Recommendations
        if recommendations:
            self._add_recommendations_slide(recommendations)
        
        # 12. Conclusion
        # 12. Conclusion
        if overall_score >= 90:
            key_message = "Excellent work! This presentation demonstrates strong academic quality across all components."
        elif overall_score >= 75:
            key_message = "Good presentation with solid fundamentals. Address the recommended improvements to reach excellence."
        elif overall_score >= 60:
            key_message = "Acceptable work with room for improvement. Focus on the critical areas identified in this report."
        elif overall_score >= 50:
            key_message = "Needs significant improvement. Review the recommendations carefully and revise your presentation."
        else:
            key_message = "This presentation requires major revision. Please address all critical issues before resubmission."
        
        self._add_conclusion_slide(grade, key_message)
        
        # Save presentation
        if output_path is None:
            output_dir = Path("outputs/reports")
            output_dir.mkdir(parents=True, exist_ok=True)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = output_dir / f"analysis_report_{timestamp}.pptx"
        else:
            output_path = Path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)
        
        self.prs.save(str(output_path))
        return str(output_path)


def generate_ppt_report(analysis_data: dict, output_path: str = None) -> str:
    """
    Convenience function to generate PowerPoint report
    
    Args:
        analysis_data: Complete analysis results dictionary
        output_path: Optional custom output path
        
    Returns:
        Path to generated PPTX file
    """
    generator = PPTReportGenerator()
    return generator.generate_report(analysis_data, output_path)
