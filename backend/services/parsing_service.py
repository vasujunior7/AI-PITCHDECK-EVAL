"""
📄 Elite Parsing Service
Extracts slides from PPTX and PDF with full metadata

Features:
- Slide titles, body text, notes
- Text box coordinates (layout regions)
- Image extraction and metadata
- Word counts and layout density
"""

from typing import List, Optional, Tuple
from pathlib import Path
import io

# PPTX parsing
from pptx import Presentation
from pptx.util import Inches, Pt

# PDF parsing
import fitz  # PyMuPDF

# Image processing
from PIL import Image
import cv2
import numpy as np

from models.slide import Slide, TextBox, ImageInfo
from core.logging import logger


def parse_slides(file_path: str) -> List[Slide]:
    """
    Parse slides from PPTX or PDF file
    
    Auto-detects format and uses appropriate parser
    
    Args:
        file_path: Path to presentation file
        
    Returns:
        List of Slide objects with complete metadata
    """
    
    file_path = Path(file_path)
    
    if not file_path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")
    
    extension = file_path.suffix.lower()
    
    logger.info(f"Parsing {extension} file: {file_path.name}")
    
    if extension == '.pptx':
        slides = parse_pptx(file_path)
    elif extension == '.pdf':
        slides = parse_pdf(file_path)
    else:
        raise ValueError(f"Unsupported file format: {extension}")
    
    logger.info(f"Successfully parsed {len(slides)} slides")
    
    return slides


def parse_pptx(file_path: Path) -> List[Slide]:
    """
    Parse PowerPoint PPTX file
    
    Extracts:
    - Slide titles (first shape with title placeholder)
    - Body text from all text boxes
    - Text box coordinates
    - Images with metadata
    - Slide notes
    """
    
    logger.info("Parsing PPTX file...")
    
    try:
        prs = Presentation(str(file_path))
    except Exception as e:
        logger.error(f"Failed to open PPTX: {str(e)}")
        raise ValueError(f"Invalid PPTX file: {str(e)}")
    
    slides = []
    
    for slide_idx, slide in enumerate(prs.slides, start=1):
        
        # Extract title
        title = extract_title_from_slide(slide)
        
        # Extract all text and text boxes
        text_boxes, body_text = extract_text_boxes(slide)
        
        # Extract notes
        notes = extract_notes(slide)
        
        # Extract images
        images = extract_images_from_pptx_slide(slide, file_path.stem, slide_idx)
        
        # Calculate metrics
        total_words = count_words(body_text)
        layout_density = calculate_layout_density(text_boxes, slide.shapes)
        
        slide_obj = Slide(
            slide_number=slide_idx,
            title=title,
            body_text=body_text,
            notes=notes,
            text_boxes=text_boxes,
            images=images,
            image_count=len(images),
            total_words=total_words,
            layout_density=layout_density
        )
        
        slides.append(slide_obj)
        
        logger.info(f"Parsed slide {slide_idx}: '{title}' ({total_words} words, {len(images)} images)")
    
    return slides


def parse_pdf(file_path: Path) -> List[Slide]:
    """
    Parse PDF file using PyMuPDF
    
    Note: PDFs don't have structured slides, so we treat each page as a slide
    """
    
    logger.info("Parsing PDF file...")
    
    try:
        doc = fitz.open(str(file_path))
    except Exception as e:
        logger.error(f"Failed to open PDF: {str(e)}")
        raise ValueError(f"Invalid PDF file: {str(e)}")
    
    slides = []
    
    for page_num in range(len(doc)):
        page = doc[page_num]
        
        # Extract text
        text = page.get_text("text")
        
        # Try to extract title (first line or largest text)
        title = extract_title_from_pdf_page(page)
        
        # Extract text blocks with positions
        text_blocks = page.get_text("dict")["blocks"]
        text_boxes = []
        body_parts = []
        
        for block in text_blocks:
            if block.get("type") == 0:  # Text block
                block_text = ""
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        block_text += span.get("text", "") + " "
                
                if block_text.strip():
                    x0, y0, x1, y1 = block["bbox"]
                    text_boxes.append(TextBox(
                        text=block_text.strip(),
                        x=float(x0),
                        y=float(y0),
                        width=float(x1 - x0),
                        height=float(y1 - y0)
                    ))
                    body_parts.append(block_text.strip())
        
        body_text = " ".join(body_parts)
        
        # Extract images
        images = extract_images_from_pdf_page(page, file_path.stem, page_num + 1)
        
        # Calculate metrics
        total_words = count_words(body_text)
        layout_density = len(text_boxes) / 100.0  # Simplified for PDF
        
        slide_obj = Slide(
            slide_number=page_num + 1,
            title=title,
            body_text=body_text,
            notes=None,  # PDFs don't have notes
            text_boxes=text_boxes,
            images=images,
            image_count=len(images),
            total_words=total_words,
            layout_density=layout_density
        )
        
        slides.append(slide_obj)
        
        logger.info(f"Parsed PDF page {page_num + 1}: '{title}' ({total_words} words, {len(images)} images)")
    
    doc.close()
    return slides


def extract_title_from_slide(slide) -> Optional[str]:
    """Extract title from PPTX slide"""
    
    # Try to find title placeholder
    if slide.shapes.title:
        title_text = slide.shapes.title.text.strip()
        if title_text:
            return title_text
    
    # Fallback: find largest text at top of slide
    texts_with_size = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            # Check if it's at the top (y < 2 inches)
            if shape.top < Inches(2):
                texts_with_size.append((shape.text.strip(), shape.top))
    
    if texts_with_size:
        # Return text at topmost position
        texts_with_size.sort(key=lambda x: x[1])
        return texts_with_size[0][0]
    
    return "Untitled Slide"


def extract_text_boxes(slide) -> Tuple[List[TextBox], str]:
    """Extract all text boxes with coordinates and combined body text"""
    
    text_boxes = []
    body_parts = []
    
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            
            # Skip if it's the title
            if slide.shapes.title and shape == slide.shapes.title:
                continue
            
            # Create text box with coordinates
            text_box = TextBox(
                text=text,
                x=float(shape.left),
                y=float(shape.top),
                width=float(shape.width),
                height=float(shape.height)
            )
            text_boxes.append(text_box)
            body_parts.append(text)
    
    body_text = " ".join(body_parts)
    
    return text_boxes, body_text


def extract_notes(slide) -> Optional[str]:
    """Extract speaker notes from slide"""
    
    try:
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            notes = text_frame.text.strip()
            return notes if notes else None
    except:
        pass
    
    return None


def extract_images_from_pptx_slide(slide, filename_base: str, slide_num: int) -> List[ImageInfo]:
    """Extract images from PPTX slide"""
    
    images = []
    image_idx = 1
    
    for shape in slide.shapes:
        # Check if shape is a picture
        if shape.shape_type == 13:  # Picture type
            try:
                image = shape.image
                
                # Get image bytes
                image_bytes = image.blob
                
                # Save image temporarily to get dimensions
                img_pil = Image.open(io.BytesIO(image_bytes))
                width, height = img_pil.size
                img_format = img_pil.format or "unknown"
                
                # Create unique path (we don't save it, just for reference)
                image_path = f"{filename_base}_slide{slide_num}_img{image_idx}.{img_format.lower()}"
                
                images.append(ImageInfo(
                    path=image_path,
                    width=width,
                    height=height,
                    format=img_format
                ))
                
                image_idx += 1
                
            except Exception as e:
                logger.warning(f"Failed to extract image from slide {slide_num}: {str(e)}")
    
    return images


def extract_images_from_pdf_page(page, filename_base: str, page_num: int) -> List[ImageInfo]:
    """Extract images from PDF page"""
    
    images = []
    image_list = page.get_images()
    
    for img_idx, img_info in enumerate(image_list, start=1):
        try:
            xref = img_info[0]
            base_image = page.parent.extract_image(xref)
            
            image_bytes = base_image["image"]
            img_ext = base_image["ext"]
            
            # Get dimensions
            img_pil = Image.open(io.BytesIO(image_bytes))
            width, height = img_pil.size
            
            image_path = f"{filename_base}_page{page_num}_img{img_idx}.{img_ext}"
            
            images.append(ImageInfo(
                path=image_path,
                width=width,
                height=height,
                format=img_ext
            ))
            
        except Exception as e:
            logger.warning(f"Failed to extract image from PDF page {page_num}: {str(e)}")
    
    return images


def extract_title_from_pdf_page(page) -> str:
    """Extract title from PDF page (heuristic: largest text at top)"""
    
    blocks = page.get_text("dict")["blocks"]
    
    title_candidates = []
    
    for block in blocks:
        if block.get("type") == 0:  # Text block
            y0 = block["bbox"][1]
            
            # Only consider text in top 20% of page
            if y0 < page.rect.height * 0.2:
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        text = span.get("text", "").strip()
                        size = span.get("size", 0)
                        
                        if text and len(text) > 3:  # Ignore very short text
                            title_candidates.append((text, size, y0))
    
    if title_candidates:
        # Sort by size (descending) then by y position (ascending)
        title_candidates.sort(key=lambda x: (-x[1], x[2]))
        return title_candidates[0][0]
    
    return "Untitled Page"


def calculate_layout_density(text_boxes: List[TextBox], shapes) -> float:
    """
    Calculate layout density (how crowded the slide is)
    
    Returns value between 0.0 (sparse) and 1.0 (very crowded)
    """
    
    if not text_boxes:
        return 0.0
    
    # Simple heuristic: number of text boxes + text length
    num_boxes = len(text_boxes)
    total_chars = sum(len(tb.text) for tb in text_boxes)
    
    # Normalize
    density = min(1.0, (num_boxes * 0.1) + (total_chars / 1000.0))
    
    return round(density, 3)


def count_words(text: str) -> int:
    """Count words in text"""
    if not text:
        return 0
    return len(text.split())
